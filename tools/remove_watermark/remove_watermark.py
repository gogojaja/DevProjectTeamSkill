#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
remove_watermark.py — 文档/图片去水印工具（v1.0.0）
支持：Word(.docx) / PPT(.pptx) / Excel(.xlsx) / PDF(.pdf) / 图片(.png/.jpg/.jpeg/.gif/.bmp/.webp) / 纯文本(.txt/.md/.log/.csv)

去水印策略（按格式）：
  - Word  ：删除页眉/页脚中的文字水印（v:textpath）与图片水印（w:pict 内 v:shape）
  - PPT   ：删除幻灯片/母版/版式中文字或图片水印形状
  - Excel ：删除工作表背景图片水印（sheetPr/picture）；--auto 额外清 header/footer 图水印
  - PDF   ：文本水印按关键字 redact；--rect 区域压制；--auto 跨页重复文字水印识别（需 PyMuPDF）
  - 图片  ：--rect 指定区域填充；--corner 常见角落水印；--auto 边缘密度启发式（PIL，无需 OpenCV）
  - 文本  ：删除整行水印印章行与行尾水印短语

可选依赖（缺失时优雅跳过并提示）：
  - PDF 需 PyMuPDF:  pip install pymupdf
  - 图片 --fill blur 用 PIL（内置）；OpenCV 可选提升质量
"""

import os
import sys
import re
import csv
import io
import shutil
import zipfile
import datetime
import argparse
import tempfile
from pathlib import Path

sys.stdout.reconfigure(encoding="utf-8")

try:
    from PIL import Image, ImageFilter, ImageStat, ImageChops
    HAVE_PIL = True
except ImportError:
    HAVE_PIL = False

try:
    import fitz  # PyMuPDF
    HAVE_FITZ = True
except ImportError:
    HAVE_FITZ = False


# ---------------------------------------------------------------------------
# 命名空间（docx/pptx zip 内 XML）
# ---------------------------------------------------------------------------
NS = {
    "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
    "v": "urn:schemas-microsoft-com:vml",
    "o": "urn:schemas-microsoft-com:office:office",
    "mc": "http://schemas.openxmlformats.org/markup-compatibility/2006",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "wp": "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing",
}
NS_RE = {k: re.escape(v) for k, v in NS.items()}
V_TEXT = f"{{{NS['v']}}}textpath"
W_P = f"{{{NS['w']}}}p"
W_PICT = f"{{{NS['w']}}}pict"
W_DRAWING = f"{{{NS['w']}}}drawing"
MC_AC = f"{{{NS['mc']}}}AlternateContent"
MC_CHOICE = f"{{{NS['mc']}}}Choice"
MC_FALLBACK = f"{{{NS['mc']}}}Fallback"
WP_ANCHOR = f"{{{NS['wp']}}}anchor"
WP_INLINE = f"{{{NS['wp']}}}inline"

LXML = None
try:
    from lxml import etree as LXML
except ImportError:
    pass


def _lxml_available():
    return LXML is not None


def make_parser():
    from lxml import etree
    parser = etree.XMLParser(remove_blank_text=False)
    return parser


def parse_xml_tolerant(xml_bytes_or_str):
    """解析 office XML。若根元素缺少 r: 等前缀声明则先补全再解析（容错真实/手工文件）。"""
    xml = xml_bytes_or_str.decode("utf-8") if isinstance(xml_bytes_or_str, bytes) else xml_bytes_or_str
    if "xmlns:r=" not in xml and "r:id" in xml:
        import re as _re
        m = _re.search(r"<([\w:]+)\s", xml)
        if m:
            root_tag = m.group(1)
            anchor = _re.search(r"<" + _re.escape(root_tag) + r"(\s[^>]*?)>", xml)
            if anchor:
                rdecl = ' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"'
                # 插到根标签第一个属性之后（保留既有属性）
                seg = anchor.group(1)
                xml = xml[:anchor.start(1)] + " " + rdecl + seg + xml[anchor.end(1):]
    return LXML.fromstring(xml.encode("utf-8"))


# ---------------------------------------------------------------------------
# zip 读取/写出工具（docx/xlsx 均走 zip，保留其余成员与压缩）
# ---------------------------------------------------------------------------

def read_zip(path):
    """读取 zip 为 {name: bytes}（保持顺序）。非 zip 返回 None。"""
    try:
        with zipfile.ZipFile(path) as z:
            return {n: z.read(n) for n in z.namelist()}
    except zipfile.BadZipFile:
        return None


def write_zip(path, members):
    """按原顺序写回 zip，保持压缩。"""
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        for name, data in members.items():
            info = zipfile.ZipInfo(name, date_time=(1980, 1, 1, 0, 0, 0))
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = 0o644 << 16 if not name.startswith("[") else 0
            z.writestr(info, data)


# ---------------------------------------------------------------------------
# Word 处理器
# ---------------------------------------------------------------------------

def watermark_xml_fragments(members):
    """定位可能含水印的部件名（headers/footers + document）。"""
    return [n for n in members
            if n.lower().startswith("word/")
            and n.lower().endswith(".xml")
            and ("header" in n.lower() or "footer" in n.lower() or "document.xml" in n.lower())]


def process_word(abspath, text, auto):
    members = read_zip(abspath)
    if members is None:
        return 0, "非有效 docx(zip)"
    if not _lxml_available():
        return 0, "需 lxml (pip install lxml)"

    total = 0
    for name in watermark_xml_fragments(members):
        xml = members[name].decode("utf-8")
        if "textpath" not in xml and "picture" not in xml:
            continue
        root = parse_xml_tolerant(xml)
        removed_ids = set()

        # 1) v:textpath 文字水印：命中关键字或 --auto
        for tp in root.findall(f".//{V_TEXT}"):
            if auto or (text and text.lower() in tp.get("string", "").lower()):
                para = tp if tp.tag == W_P else tp.getparent()
                while para is not None and para.tag != W_P:
                    para = para.getparent()
                if para is not None:
                    removed_ids.add(id(para))

        # 2) w:pict 内 v:shape 图片水印（PowerPlusWaterMarkObject / 水印命名）
        for pict in root.findall(f".//{W_PICT}"):
            shape = pict.find(f".//{{{NS['v']}}}shape")
            if shape is None:
                continue
            sid = shape.get("id") or ""
            sname = shape.get("name") or ""
            if ("Watermark" in sid or "水印" in sname or "Watermark" in sname) and (auto or not text):
                para = pict if pict.tag == W_P else pict.getparent()
                while para is not None and para.tag != W_P:
                    para = para.getparent()
                if para is not None:
                    removed_ids.add(id(para))

        if not removed_ids:
            continue
        cnt = 0
        for para in list(root.iter(W_P)):
            if id(para) in removed_ids:
                para.getparent().remove(para)
                cnt += 1
        if cnt:
            members[name] = LXML.tostring(root, xml_declaration=True,
                                          encoding="UTF-8", standalone=True)
            total += cnt
    if total:
        write_zip(abspath, members)
    return total, f"删除 {total} 个水印段落" if total else "未发现文字水印"


def _sync_rels(members, *_):
    """占位：word 水印为页眉段落内图表，无需改 rels（表头部件仍被引用）。"""
    return members


# ---------------------------------------------------------------------------
# PPT 处理器
# ---------------------------------------------------------------------------

def process_ppt(abspath, text, auto):
    try:
        from pptx import Presentation
        from pptx.util import Length
    except ImportError:
        return 0, "需 python-pptx (pip install python-pptx)"

    prs = Presentation(abspath)
    cnt = 0

    def _match(sh):
        try:
            if sh.has_text_frame:
                t = sh.text_frame.text or ""
            else:
                t = ""
        except Exception:
            t = ""
        name = (sh.name or "") if hasattr(sh, "name") else ""
        if auto:
            # 自动：大字号/旋转 文字 或 名称含 Watermark/水印
            if "watermark" in name.lower() or "水印" in name:
                return True
            if sh.has_text_frame and sh.text_frame.text and len(sh.text_frame.text) < 30:
                try:
                    size = None
                    for para in sh.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                size = run.font.size
                    if size and size >= Length(400000):  # >= 32pt
                        return True
                except Exception:
                    pass
        return bool(text and text.lower() in t.lower())

    def _scan(shapes):
        nonlocal cnt
        for sh in shapes:
            if _match(sh):
                el = sh._element
                el.getparent().remove(el)
                cnt += 1
            if sh.shape_type is not None and sh.shape_type == 6:  # GROUP
                _scan(sh.shapes)

    for slide in prs.slides:
        _scan(slide.shapes)
    for master in prs.slide_masters:
        _scan(master.shapes)
    for layout in prs.slide_layouts:
        _scan(layout.shapes)

    if cnt:
        prs.save(abspath)
    return cnt, f"删除 {cnt} 个水印形状" if cnt else "未发现水印形状"


# ---------------------------------------------------------------------------
# Excel 处理器（zip 级：删除工作稿背景 picture）
# ---------------------------------------------------------------------------

def _remove_xlsx_background(members, auto):
    removed_sheets = 0
    removed_hf = 0
    # 1) 各工作表 sheetN.xml 中 sheetPr/picture 背景图
    for name in [n for n in members if re.match(r"xl/worksheets/sheet\d+\.xml$", n)]:
        xml = members[name].decode("utf-8")
        if "<picture" not in xml:
            continue
        if not _lxml_available():
            continue
        root = parse_xml_tolerant(xml)
        sheet_pr = root.find(".//{http://schemas.openxmlformats.org/spreadsheetml/2006/main}sheetPr")
        if sheet_pr is None:
            continue
        pic = sheet_pr.find("{http://schemas.openxmlformats.org/spreadsheetml/2006/main}picture")
        if pic is not None:
            sheet_pr.remove(pic)
            removed_sheets += 1
        members[name] = LXML.tostring(root, xml_declaration=True,
                                      encoding="UTF-8", standalone=True)

    # 2) --auto：header/footer 图片水印（legacyDrawingHF → vmlDrawing + media 图）
    if auto:
        for name in [n for n in members if re.match(r"xl/worksheets/sheet\d+\.xml$", n)]:
            xml = members[name].decode("utf-8")
            if "legacyDrawingHF" not in xml:
                continue
            if not _lxml_available():
                continue
            root = parse_xml_tolerant(xml)
            # excel: <legacyDrawingHF r:id="rIdX"/> 默认命名空间
            default_ns = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
            for el in root.iter("{%s}legacyDrawingHF" % default_ns):
                el.getparent().remove(el)
                removed_hf += 1
            members[name] = LXML.tostring(root, xml_declaration=True,
                                          encoding="UTF-8", standalone=True)
    return removed_sheets, removed_hf


def process_excel(abspath, text, auto):
    members = read_zip(abspath)
    if members is None:
        return 0, "非有效 xlsx(zip)"
    if not _lxml_available():
        return 0, "需 lxml (pip install lxml)"
    sheets, hf = _remove_xlsx_background(members, auto)
    if sheets or hf:
        write_zip(abspath, members)
    msg = []
    if sheets:
        msg.append(f"清除 {sheets} 个工作表背景水印")
    if hf:
        msg.append(f"清除 {hf} 处页眉/页脚图水印")
    return sheets + hf, "; ".join(msg) or "未发现背景/页眉水印"


# ---------------------------------------------------------------------------
# PDF 处理器（PyMuPDF 可选）
# ---------------------------------------------------------------------------

def _pdf_text_spans(page):
    """返回 [(bbox, text, flags)] 文本片段。"""
    out = []
    d = page.get_text("dict")
    for block in d.get("blocks", []):
        for line in block.get("lines", []):
            for span in line.get("spans", []):
                if span.get("text", "").strip():
                    out.append((span["bbox"], span["text"], span.get("flags", 0)))
    return out


def process_pdf(abspath, text, auto, rects):
    if not HAVE_FITZ:
        return 0, "需 PyMuPDF（pip install pymupdf）——跳过 PDF"
    doc = fitz.open(abspath)
    total = 0
    used_rects = [list(r) for r in rects]

    # 跨页重复文字水印识别（--auto）：相同 bbox 尺寸+极高频率的短文本
    candidates = {}
    if auto:
        for pno in range(doc.page_count):
            for bbox, t, flags in _pdf_text_spans(doc[pno]):
                key = (round(bbox[2] - bbox[0], 1), round(bbox[3] - bbox[1], 1), t)
                candidates[key] = candidates.get(key, 0) + 1
        threshold = max(3, int(doc.page_count * 0.5))
        auto_keys = {k for k, c in candidates.items() if c >= threshold}

    for pno in range(doc.page_count):
        page = doc[pno]
        hits = []

        # 文本关键字命中
        if text:
            for bbox, t, flags in _pdf_text_spans(page):
                if text.lower() in t.lower():
                    hits.append(fitz.Rect(bbox))

        # 跨页重复（auto）
        if auto:
            for bbox, t, flags in _pdf_text_spans(page):
                key = (round(bbox[2] - bbox[0], 1), round(bbox[3] - bbox[1], 1), t)
                if key in auto_keys:
                    hits.append(fitz.Rect(bbox))

        # --rect 区域压制
        for r in used_rects:
            x0, y0, x1, y1 = r
            if x1 > x0 and y1 > y0:
                hits.append(fitz.Rect(x0 * page.rect.width, y0 * page.rect.height,
                                      x1 * page.rect.width, y1 * page.rect.height))
            else:
                hits.append(fitz.Rect(x0, y0, max(x0, x1), max(y0, y1)))

        if hits:
            for r in hits:
                page.add_redact_annot(r, fill=(1, 1, 1))
            page.apply_redactions()
            total += len(hits)

    doc.save(abspath, garbage=4, deflate=True)
    doc.close()
    return total, f"压制 {total} 处水印" if total else "未发现水印"


# ---------------------------------------------------------------------------
# 图片处理器（PIL）
# ---------------------------------------------------------------------------

def image_corner_rect(w, h, corner, frac=0.14):
    fw, fh = max(int(w * frac), 40), max(int(h * frac), 40)
    if corner == "br":
        return (w - fw, h - fh, w, h)
    if corner == "tr":
        return (w - fw, 0, w, fh)
    if corner == "bl":
        return (0, h - fh, fw, h)
    return (0, 0, fw, fh)  # tl


def _edge_density(im):
    g = im.convert("L").filter(ImageFilter.FIND_EDGES)
    return ImageStat.Stat(g).mean[0]


def process_image(abspath, text, auto, rects, corner, fill):
    if not HAVE_PIL:
        return 0, "需 Pillow（pip install pillow）"
    im = Image.open(abspath)
    im.load()
    W, H = im.size

    # 1) 显式 --rect
    fill_rects = []
    for r in rects:
        if len(r) == 4:
            x0, y0, x1, y1 = r
            fill_rects.append((min(x0, x1), min(y0, y1), max(x0, x1), max(y0, y1)))
    # 2) --corner
    if corner:
        fill_rects.append(image_corner_rect(W, H, corner))
    # 3) --auto 边缘密度启发式：对比四角，取"最像水印"即边缘密度最低且位于角落的文字带
    if auto and not fill_rects:
        scores = {}
        for c in ("tl", "tr", "bl", "br"):
            x0, y0, x1, y1 = image_corner_rect(W, H, c, 0.12)
            crop = im.crop((x0, y0, x1, y1))
            den = _edge_density(crop)
            scores[c] = den
        # 水印文字稀疏 → 边缘密度显著低于其他角（带文字的角应高于背景，故取最高）
        best = max(scores, key=scores.get)
        # 仅当该角平均亮度偏中灰偏移（非纯背景内嵌）——保守：取密度最大的角
        fill_rects.append(image_corner_rect(W, H, best, 0.12))

    if not fill_rects:
        return 0, "未指定水印区域（用 --rect 或 --corner 或 --auto）"

    if fill == "blur":
        base = im.filter(ImageFilter.GaussianBlur(radius=8))
    elif fill == "edge":
        base = _edge_fill(im)
    else:  # white
        base = Image.new("RGB", im.size, (255, 255, 255))

    for r in fill_rects:
        x0, y0, x1, y1 = r
        im.paste(base.crop((x0, y0, x1, y1)), (x0, y0))

    if im.mode in ("RGBA", "LA") or (im.mode == "P" and "transparency" in im.info):
        im.save(abspath)
    else:
        im = im.convert("RGB")
        im.save(abspath)
    return len(fill_rects), f"填充 {len(fill_rects)} 个区域"


def _edge_fill(im):
    """按四边平均色生成的纯色底（用于水印区域覆盖）。"""
    W, H = im.size
    edges = [im.crop((0, 0, W, 2)), im.crop((0, H - 2, W, H)),
             im.crop((0, 0, 2, H)), im.crop((W - 2, 0, W, H))]
    means = [tuple(int(x) for x in ImageStat.Stat(e).mean) for e in edges]
    avg = tuple(sum(c[i] for c in means) // len(means) for i in range(3))
    return Image.new("RGB", im.size, avg)


# ---------------------------------------------------------------------------
# 文本处理器
# ---------------------------------------------------------------------------

def process_text(abspath, text, auto, encoding="utf-8"):
    try:
        content = Path(abspath).read_text(encoding=encoding)
    except UnicodeDecodeError:
        for enc in ("utf-8-sig", "gb18030", "utf-16"):
            try:
                content = Path(abspath).read_text(encoding=enc)
                encoding = enc
                break
            except (UnicodeDecodeError, LookupError):
                continue
        else:
            return 0, "编码无法识别"
    lines = content.splitlines(keepends=True)
    cnt = 0

    def _stamp(line):
        # 整行印章：行 == 水印文本，或行 = 水印文本 + 少量尾注（如 "内部资料 勿外传"）
        s = line.strip("\r\n").strip()
        if not (text and s):
            return False
        if s == text.strip():
            return True
        if s.startswith(text.strip()):
            tail = s[len(text.strip()):].strip()
            return bool(tail) and len(tail) <= 4 and not any(c.isascii() and c.isalnum() for c in tail)
        return False

    out = []
    for line in lines:
        if _stamp(line):
            cnt += 1
            continue
        # --text：行尾水印短语（保持非印章行主体）
        if text and not auto:
            s = line
            stripped = s.rstrip("\r\n")
            trail = stripped[len(stripped.rstrip()):]
            body = stripped.rstrip()
            if body.endswith(text):
                line = body[: -len(text)] + trail
                cnt += 1
        # --auto：清洗高频整行印章（长度 ≤24 且完全相同的行出现 ≥3 次）
        out.append(line)

    # auto：统计并删除高频重复行（印章带）
    if auto:
        from collections import Counter
        body_lines = [l.strip("\r\n") for l in out]
        counts = Counter(body_lines)
        stamps = {b for b, c in counts.items() if c >= 3 and 1 <= len(b) <= 24}
        if stamps:
            new_out = []
            for line in out:
                if line.strip("\r\n") in stamps:
                    cnt += 1
                    continue
                new_out.append(line)
            out = new_out

    if cnt:
        Path(abspath).write_text("".join(out), encoding=encoding)
    return cnt, f"删除 {cnt} 处水印行/短语" if cnt else "未发现水印行"


# ---------------------------------------------------------------------------
# 分发
# ---------------------------------------------------------------------------

TEXT_EXTS = {".txt", ".md", ".log", ".csv", ".rst", ".adoc", ".xml", ".yaml", ".yml", ".json", ".ini", ".conf"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tif", ".tiff"}


def dispatch_kind(path, forced):
    ext = Path(path).suffix.lower()
    if forced:
        return forced
    if ext in (".docx",):
        return "word"
    if ext in (".pptx", ".ppt"):
        return "ppt"
    if ext in (".xlsx", ".xlsm"):
        return "excel"
    if ext == ".pdf":
        return "pdf"
    if ext in IMAGE_EXTS:
        return "image"
    if ext in TEXT_EXTS:
        return "text"
    return "text"  # 未知扩展名按文本处理


def process_file(fp, text, auto, rects, corner, fill, forced, kind=None):
    kind = kind or dispatch_kind(fp, forced)
    if kind == "word":
        return process_word(fp, text, auto)
    if kind == "ppt":
        return process_ppt(fp, text, auto)
    if kind == "excel":
        return process_excel(fp, text, auto)
    if kind == "pdf":
        return process_pdf(fp, text, auto, rects)
    if kind == "image":
        return process_image(fp, text, auto, rects, corner, fill)
    return process_text(fp, text, auto)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="去水印工具 — Word/PPT/Excel/PDF/图片/纯文本（按格式自动识别处理器）",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  # Word：打开并删除页眉文字水印（--auto 识别；--text 指定文案）
  python remove_watermark.py report.docx --in-place --auto
  python remove_watermark.py report.docx --text "内部资料"

  # PPT：删除含指定文字的水印形状
  python remove_watermark.py deck.pptx --text "机密"

  # Excel：清除工作表背景图水印（--auto 额外清页眉页脚图水印）
  python remove_watermark.py book.xlsx --in-place
  python remove_watermark.py book.xlsx --in-place --auto

  # PDF：文字水印按关键字 redact；--rect 相对坐标(0-1)区域压制
  python remove_watermark.py doc.pdf --text "CONFIDENTIAL"
  python remove_watermark.py doc.pdf --rect 0.05,0.85,0.95,0.95

  # 图片：--corner 角落 / --rect 像素区域 / --auto 启发式
  python remove_watermark.py photo.png --corner br --fill edge -o out/
  python remove_watermark.py photo.png --rect 300,200,800,500
  python remove_watermark.py shot.jpg --auto --fill blur

  # 文本：删除水印印章行/行尾短语
  python remove_watermark.py notes.md --text "内部资料 勿外传"
  python remove_watermark.py log.txt --auto
""")
    parser.add_argument("target", nargs="?", help="目标文件或目录")
    parser.add_argument("-o", "--output", help="输出目录（默认同目录 _nowater 导出，保留原件）")
    parser.add_argument("--in-place", action="store_true", help="原地修改（危险，先备份；--auto 配合）")
    parser.add_argument("--text", help="水印文字/文案（关键字匹配）")
    parser.add_argument("--auto", action="store_true", help="自动识别水印（各格式启发式）")
    parser.add_argument("--rect", action="append", default=[], metavar="x0,y0,x1,y1",
                        help="区域（PDF 用相对 0-1；图片用像素）。可多次——包围多个水印")
    parser.add_argument("--corner", choices=["tl", "tr", "bl", "br"], help="图片：常见角落水印")
    parser.add_argument("--fill", choices=["blur", "edge", "white"], default="blur", help="图片填充方式（默认 blur 模糊化）")
    parser.add_argument("--format", choices=["word", "ppt", "excel", "pdf", "image", "text"], help="强制按指定格式处理")
    parser.add_argument("--dry-run", action="store_true", help="预览：仅报告将处理的文件，不修改")
    parser.add_argument("--report", help="报告 CSV 路径")
    parser.add_argument("--include-ext", help="额外扩展名（文本，逗号分隔，如 .conf,.properties）")
    args = parser.parse_args()

    if not args.target:
        parser.error("请指定目标文件或目录")

    if args.in_place and args.output:
        parser.error("--in-place 与 -o 互斥")

    target = Path(args.target)
    if not target.exists():
        parser.error(f"目标不存在: {args.target}")

    rects = []
    for rstr in args.rect:
        parts = [p.strip() for p in rstr.split(",")]
        if len(parts) != 4:
            parser.error(f"--rect 需 4 个数字: {rstr}")
        rects.append(tuple(float(p) for p in parts))

    # 文件收集
    text_exts = set(TEXT_EXTS)
    if args.include_ext:
        for e in args.include_ext.split(","):
            e = e.strip().lower()
            if e and not e.startswith("."):
                e = "." + e
            text_exts.add(e)

    files = []
    if target.is_file():
        files = [target]
    else:
        for root, dirs, names in os.walk(target):
            dirs[:] = [d for d in dirs if d not in {".git", ".venv", "__pycache__", "dist", "build", ".secrets"}]
            for n in names:
                files.append(Path(root) / n)

    results = []
    for fp in files:
        kind = dispatch_kind(str(fp), args.format)
        if kind == "text" and not args.format and fp.suffix.lower() not in text_exts and fp.suffix.lower() not in {".txt", ".md", ".log", ".csv"}:
            continue  # 目录扫描时非文本扩展名跳过，除非强制 text
        if kind not in {"word", "ppt", "excel", "pdf", "image", "text"}:
            continue
        if args.dry_run:
            results.append({"file": str(fp), "kind": kind, "status": "DRY_RUN", "detail": "预览不修改"})
            continue

        # 输出策略：原地 or 复制到 -o
        if args.in_place:
            work = str(fp)
        else:
            out_dir = Path(args.output) if args.output else fp.parent / "_nowater"
            out_path = out_dir / fp.name if not target.is_dir() else out_dir / fp.relative_to(target)
            out_path.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(fp, out_path)
            work = str(out_path)

        try:
            n, detail = process_file(work, args.text, args.auto, rects, args.corner, args.fill, args.format, kind)
            if n:
                status = "OK"
            else:
                status = "NO_WATERMARK"
            results.append({"file": str(fp), "kind": kind, "status": status, "detail": detail, "count": n})
            print(f"  {'✅' if n else '·'} {fp.name} ({kind}) — {detail}")
        except Exception as e:
            results.append({"file": str(fp), "kind": kind, "status": "ERROR", "detail": str(e)})
            print(f"  ❌ {fp.name} ({kind}) — 失败: {e}")

    # 报告
    if args.report:
        with open(args.report, "w", newline="", encoding="utf-8-sig") as f:
            w = csv.writer(f)
            w.writerow(["file", "kind", "status", "detail"])
            for r in results:
                w.writerow([r["file"], r.get("kind"), r.get("status"), r.get("detail", "")])
        print(f"\n📊 报告: {args.report}")

    ok = sum(1 for r in results if r.get("status") == "OK")
    nowm = sum(1 for r in results if r.get("status") == "NO_WATERMARK")
    err = sum(1 for r in results if r.get("status") == "ERROR")
    print(f"\n处理 {len(results)} 个文件：清理 {ok}，无水印 {nowm}，失败 {err}")
    return 1 if err else 0


if __name__ == "__main__":
    sys.exit(main())